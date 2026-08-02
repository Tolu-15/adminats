#!/usr/bin/env python3
"""
Generate a professional PowerPoint presentation for the ATS Membership App.
Covers: overview, architecture, workflow, features, tech stack, and future.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Colour palette ──────────────────────────────────────────────
NAVY       = RGBColor(0x0D, 0x1B, 0x2E)
DARK_NAVY  = RGBColor(0x07, 0x0E, 0x1A)
GOLD       = RGBColor(0xD4, 0xAF, 0x37)
BLUE       = RGBColor(0x3B, 0x82, 0xF6)
PURPLE     = RGBColor(0x8B, 0x5C, 0xF6)
GREEN      = RGBColor(0x10, 0xB9, 0x81)
RED        = RGBColor(0xEF, 0x44, 0x44)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF1, 0xF5, 0xF9)
MID_GRAY   = RGBColor(0x94, 0xA3, 0xB8)
AMBER      = RGBColor(0xF5, 0x9E, 0x0B)


def set_slide_bg(slide, color):
    """Set a solid background on a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text,
                 font_size=18, bold=False, color=WHITE,
                 alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_slide_content(slide, items, left, top, width, height,
                              font_size=15, color=WHITE, spacing=Pt(8)):
    """Add multiple bullet-point paragraphs to one text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
        p.level = 0
    return txBox


def add_accent_bar(slide, left, top, width, height, color=GOLD):
    """Add a small coloured rectangle accent bar."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_card(slide, left, top, width, height, bg_color=RGBColor(0x12, 0x24, 0x3B)):
    """Add a rounded-rectangle card shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    return shape


def add_section_header(slide, title, subtitle=None):
    """Standard header block for content slides."""
    add_accent_bar(slide, Inches(0.5), Inches(0.45), Inches(0.35), Inches(0.06))
    add_text_box(slide, Inches(0.5), Inches(0.6), Inches(8.5), Inches(0.55),
                 title, font_size=28, bold=True, color=WHITE)
    if subtitle:
        add_text_box(slide, Inches(0.5), Inches(1.1), Inches(8.5), Inches(0.4),
                     subtitle, font_size=14, color=MID_GRAY)


# ═════════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ═════════════════════════════════════════════════════════════════

def slide_title(prs):
    """Slide 1 — Title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, DARK_NAVY)

    # Decorative gold line
    add_accent_bar(slide, Inches(3.2), Inches(1.8), Inches(3.6), Inches(0.05))

    add_text_box(slide, Inches(0.8), Inches(2.0), Inches(8.4), Inches(1.0),
                 "ATS Membership App", font_size=40, bold=True, color=GOLD,
                 alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1.0), Inches(2.9), Inches(8.0), Inches(0.6),
                 "Workflow, Features & System Architecture",
                 font_size=20, color=WHITE, alignment=PP_ALIGN.CENTER)

    # Decorative gold line
    add_accent_bar(slide, Inches(3.2), Inches(3.55), Inches(3.6), Inches(0.05))

    add_text_box(slide, Inches(1.0), Inches(4.0), Inches(8.0), Inches(0.5),
                 "Citadel Global Community Church  •  Apostolic Training School",
                 font_size=13, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1.0), Inches(4.6), Inches(8.0), Inches(0.4),
                 "August 2026",
                 font_size=12, color=MID_GRAY, alignment=PP_ALIGN.CENTER)


def slide_agenda(prs):
    """Slide 2 — Agenda / Table of Contents."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    add_section_header(slide, "Agenda")

    items = [
        "1.  App Overview & Purpose",
        "2.  System Architecture & Tech Stack",
        "3.  Core Workflow — Registration to Grading",
        "4.  Three-Programme Model (Membership · MIT · Proclaimers)",
        "5.  Admin Dashboard & Analytics",
        "6.  Batch Management & QR Code System",
        "7.  Student Profiles & Grade Management",
        "8.  Excel Import / Export & Data Migration",
        "9.  Retake System",
        "10. Google Sheets Sync",
        "11. Security & Authentication",
        "12. Deployment & Next Steps",
    ]
    add_bullet_slide_content(slide, items,
                              Inches(0.8), Inches(1.6), Inches(8.4), Inches(4.0),
                              font_size=14, color=LIGHT_GRAY, spacing=Pt(6))


def slide_overview(prs):
    """Slide 3 — App Overview."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    add_section_header(slide, "App Overview",
                       "What is the ATS Membership App and why was it built?")

    items = [
        "▸  Replaces legacy Microsoft Forms for member registration",
        "▸  Full-stack web application for the Apostolic Training School (ATS)",
        "▸  Covers three training programmes: Membership, MIT, and Proclaimers",
        "▸  Each student receives a unique, auto-generated ID (e.g. ATS-056-0001)",
        "▸  Real-time sync with Google Sheets as a secondary record",
        "▸  Admins manage batches, view analytics, and grade students — all in one place",
        "▸  Supabase is the single source of truth for all data",
    ]
    add_bullet_slide_content(slide, items,
                              Inches(0.6), Inches(1.7), Inches(8.8), Inches(3.8),
                              font_size=15, color=WHITE, spacing=Pt(10))


def slide_architecture(prs):
    """Slide 4 — Architecture & Tech Stack."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    add_section_header(slide, "System Architecture & Tech Stack")

    # Cards for each technology
    card_data = [
        ("Frontend", "Next.js 14 + React 18", "Server-side rendered pages,\nApp Router, dynamic routes", BLUE),
        ("Backend", "Next.js API Routes", "RESTful endpoints for batches,\nstudents, grades, import/export", PURPLE),
        ("Database", "Supabase (PostgreSQL)", "RLS policies, auto-IDs,\nreal-time storage bucket", GREEN),
        ("Integrations", "Google Sheets API + XLSX", "Bidirectional sync, Excel\nimport/export, data migration", AMBER),
    ]

    for i, (title, tech, desc, accent) in enumerate(card_data):
        col = i % 4
        x = Inches(0.4 + col * 2.4)
        y = Inches(1.8)

        card = add_card(slide, x, y, Inches(2.2), Inches(2.6))

        # accent bar on card
        add_accent_bar(slide, x + Inches(0.15), y + Inches(0.15),
                       Inches(0.5), Inches(0.05), accent)

        add_text_box(slide, x + Inches(0.15), y + Inches(0.3),
                     Inches(1.9), Inches(0.35),
                     title, font_size=13, bold=True, color=accent)

        add_text_box(slide, x + Inches(0.15), y + Inches(0.7),
                     Inches(1.9), Inches(0.45),
                     tech, font_size=11, bold=True, color=WHITE)

        add_text_box(slide, x + Inches(0.15), y + Inches(1.2),
                     Inches(1.9), Inches(1.2),
                     desc, font_size=10, color=MID_GRAY)


def slide_workflow(prs):
    """Slide 5 — Core Workflow."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    add_section_header(slide, "Core Workflow",
                       "From batch creation to student grading — the end-to-end flow")

    steps = [
        ("1", "Create Batch", "Admin creates a new batch\n(code, name, programme type)\nand gets a unique reg link.", BLUE),
        ("2", "Share Link", "Link or QR code is shared\nwith students via WhatsApp\nor printed materials.", GOLD),
        ("3", "Student Registers", "Student fills in biodata,\nuploads a photo, and submits.\nUnique ID is auto-generated.", GREEN),
        ("4", "Data Stored", "Record saved to Supabase\nand appended to Google\nSheets automatically.", PURPLE),
        ("5", "Admin Reviews", "Admin views profiles,\nassigns grades, exports\ndata, or imports via Excel.", AMBER),
    ]

    for i, (num, title, desc, color) in enumerate(steps):
        x = Inches(0.25 + i * 1.93)
        y = Inches(1.9)

        # Circle for step number
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.6), y, Inches(0.5), Inches(0.5))
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()

        add_text_box(slide, x + Inches(0.6), y + Inches(0.07),
                     Inches(0.5), Inches(0.4),
                     num, font_size=18, bold=True, color=WHITE,
                     alignment=PP_ALIGN.CENTER)

        # Connector arrow (except last)
        if i < len(steps) - 1:
            add_text_box(slide, x + Inches(1.25), y + Inches(0.1),
                         Inches(0.5), Inches(0.35),
                         "→", font_size=20, bold=True, color=MID_GRAY,
                         alignment=PP_ALIGN.CENTER)

        # Step title
        add_text_box(slide, x, y + Inches(0.65),
                     Inches(1.7), Inches(0.35),
                     title, font_size=13, bold=True, color=color,
                     alignment=PP_ALIGN.CENTER)

        # Step description
        add_text_box(slide, x, y + Inches(1.05),
                     Inches(1.7), Inches(1.6),
                     desc, font_size=10, color=LIGHT_GRAY,
                     alignment=PP_ALIGN.CENTER)


def slide_three_programmes(prs):
    """Slide 6 — Three-Programme Model."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    add_section_header(slide, "Three-Programme Model",
                       "Each batch belongs to one of three ATS training levels")

    progs = [
        ("MEMBERSHIP (MEM-100)", BLUE,
         "▸ New members register via unique link\n"
         "▸ Personal & spiritual biodata collected\n"
         "▸ Photo upload + unique Student ID\n"
         "▸ Grades: Attendance, Test, Assignment,\n"
         "   Assessment, Presentation, Exam\n"
         "▸ Tracks water & Holy Spirit baptism,\n"
         "   covenant deed signing"),

        ("MIT (MIT-200)", GOLD,
         "▸ Ministry-in-Training programme\n"
         "▸ Links existing Membership students\n"
         "▸ Students look up their ID or card number\n"
         "▸ Grades: Midterm, Interactions, Bible Study,\n"
         "   CTH, Community Service, Evangelism\n"
         "▸ Tracks department + first-timer status"),

        ("PROCLAIMERS (PROC-300)", PURPLE,
         "▸ Advanced ministry programme\n"
         "▸ Also links existing Membership students\n"
         "▸ Grades: CIH, Attendance, Assessment,\n"
         "   Presentation, Project, Mountain of\n"
         "   Influence, Seminar Attendance\n"
         "▸ Department assignment per student"),
    ]

    for i, (title, color, desc) in enumerate(progs):
        x = Inches(0.35 + i * 3.15)
        y = Inches(1.7)
        card = add_card(slide, x, y, Inches(3.0), Inches(3.5))

        add_accent_bar(slide, x + Inches(0.2), y + Inches(0.2),
                       Inches(2.6), Inches(0.05), color)

        add_text_box(slide, x + Inches(0.2), y + Inches(0.35),
                     Inches(2.6), Inches(0.35),
                     title, font_size=14, bold=True, color=color)

        add_text_box(slide, x + Inches(0.2), y + Inches(0.8),
                     Inches(2.6), Inches(2.5),
                     desc, font_size=10.5, color=LIGHT_GRAY)


def slide_admin_dashboard(prs):
    """Slide 7 — Admin Dashboard & Analytics."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    add_section_header(slide, "Admin Dashboard & Analytics",
                       "Real-time statistics and visual insights at a glance")

    stat_cards = [
        ("👥", "Total Enrolled", "Membership + MIT + Proclaimers\nacross all batches", BLUE),
        ("⚤", "Gender Ratio", "Male vs Female percentage\nbreakdown with counts", PURPLE),
        ("⭐", "First Timers", "Count and percentage of\nfirst-time members", GOLD),
        ("📊", "Batches Overview", "Active vs total batches\nregistered in the system", GREEN),
    ]

    for i, (icon, label, desc, color) in enumerate(stat_cards):
        x = Inches(0.4 + i * 2.4)
        y = Inches(1.7)
        card = add_card(slide, x, y, Inches(2.2), Inches(1.4))
        add_text_box(slide, x + Inches(0.15), y + Inches(0.1),
                     Inches(0.4), Inches(0.4),
                     icon, font_size=20, color=WHITE)
        add_text_box(slide, x + Inches(0.55), y + Inches(0.15),
                     Inches(1.5), Inches(0.3),
                     label, font_size=12, bold=True, color=color)
        add_text_box(slide, x + Inches(0.15), y + Inches(0.55),
                     Inches(1.9), Inches(0.8),
                     desc, font_size=9.5, color=MID_GRAY)

    # Chart descriptions
    charts_text = [
        "📈 Batch Student Growth — Multi-line trend chart showing Membership, MIT, and Proclaimers student counts per batch with interactive hover tooltips.",
        "📊 Student Demographics — Bar chart comparing Male/Female counts and First-Timer ratios with clean, color-coded visual labels.",
    ]
    add_bullet_slide_content(slide, charts_text,
                              Inches(0.6), Inches(3.4), Inches(8.8), Inches(1.8),
                              font_size=12, color=LIGHT_GRAY, spacing=Pt(14))


def slide_batch_management(prs):
    """Slide 8 — Batch Management & QR Codes."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    add_section_header(slide, "Batch Management & QR Codes",
                       "Create, search, share, and manage registration batches")

    items = [
        "▸  Create Batch — Assign code, name, and programme type (Membership / MIT / Proclaimers)",
        "▸  Unique Registration Link — Auto-generated reg_token for each batch, shareable via URL",
        "▸  QR Code Generation — One-click QR codes via qrserver.com API, downloadable as PNG",
        "▸  Copy Link — Quick clipboard copy of the full registration URL",
        "▸  Batch Search & Filter — Search by code or name, filter by programme type",
        "▸  Pagination — 3 batch cards per page with Previous / Next navigation",
        "▸  Delete Batch — Confirmation-gated batch deletion with cascade protection",
        "▸  Batch Detail View — See all students in a batch, with programme-specific student tables",
        "▸  Activate / Deactivate — Toggle is_active flag to open or close registrations",
    ]
    add_bullet_slide_content(slide, items,
                              Inches(0.6), Inches(1.7), Inches(8.8), Inches(3.8),
                              font_size=13, color=WHITE, spacing=Pt(8))


def slide_student_profiles(prs):
    """Slide 9 — Student Profiles & Grade Management."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    add_section_header(slide, "Student Profiles & Grade Management",
                       "Comprehensive student data and inline grade editing")

    # Two column layout
    left_items = [
        "Student Profile Card:",
        "▸  Auto-generated Student ID",
        "▸  Photo (uploaded to Supabase Storage)",
        "▸  Full biodata: name, email, phone,",
        "    DOB, gender, address, next of kin",
        "▸  Spiritual info: born again status,",
        "    water/Holy Spirit baptism, church join date",
        "▸  Education level & first-timer flag",
        "▸  Card number for physical membership card",
    ]

    right_items = [
        "Grade Management:",
        "▸  Inline grade editing per student",
        "▸  Membership: 6 scores + status + comments",
        "▸  MIT: 10+ scores including CTH, evangelism",
        "▸  Proclaimers: CIH, assessment, project",
        "▸  Pass / Fail status tracking",
        "▸  Retake badge detection (🔄 icon)",
        "▸  Global student search across all batches",
        "▸  Search by name, student ID, or card no.",
    ]

    add_bullet_slide_content(slide, left_items,
                              Inches(0.4), Inches(1.7), Inches(4.5), Inches(3.8),
                              font_size=12, color=LIGHT_GRAY, spacing=Pt(6))
    add_bullet_slide_content(slide, right_items,
                              Inches(5.0), Inches(1.7), Inches(4.5), Inches(3.8),
                              font_size=12, color=LIGHT_GRAY, spacing=Pt(6))


def slide_import_export(prs):
    """Slide 10 — Excel Import / Export."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    add_section_header(slide, "Excel Import / Export & Data Migration",
                       "Bulk data operations with intelligent matching and retake detection")

    items = [
        "EXPORT (Download as .xlsx):",
        "   ▸  One-click download per batch — generates multi-sheet Excel workbook",
        "   ▸  Separate tabs: Membership, MIT, Proclaimers (only if data exists)",
        "   ▸  Student name, ID, class, trainer, all scores, status, comments",
        "",
        "IMPORT (Upload .xlsx):",
        "   ▸  Intelligent header detection — scans first 10 rows for known column names",
        "   ▸  Master Bio Sheet — creates students with auto-generated unique IDs",
        "   ▸  Grades Sheet — upserts grade records per student",
        "   ▸  Cross-batch student matching by email, card number, or name",
        "   ▸  Retake detection — automatically flags students re-enrolling from prior batches",
        "   ▸  Collision-safe ID generation with up to 5 retry attempts",
        "",
        "MIGRATION TEMPLATE:",
        "   ▸  Downloadable template Excel for standardised bulk imports",
    ]
    add_bullet_slide_content(slide, items,
                              Inches(0.6), Inches(1.6), Inches(8.8), Inches(4.0),
                              font_size=11.5, color=WHITE, spacing=Pt(4))


def slide_retake(prs):
    """Slide 11 — Retake System."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    add_section_header(slide, "Retake System",
                       "Supporting students who need to re-enrol across batches")

    retake_data = [
        ("Membership Retake", BLUE, [
            "▸ Student searches by existing ID or card number",
            "▸ API updates batch_id to the new batch",
            "▸ Grade record is RESET with 'RETAKE —' comment",
            "▸ Prior status is preserved in comments",
        ]),
        ("MIT Retake", GOLD, [
            "▸ Lookup returns isRetake=true + prior attempts",
            "▸ New mit_registrations row per batch allowed",
            "▸ New mit_grades row with 'RETAKE —' comment",
            "▸ Multiple enrolments across batches supported",
        ]),
        ("Proclaimers Retake", PURPLE, [
            "▸ Same approach as MIT retake system",
            "▸ Multiple proclaimers_registrations per student",
            "▸ Each batch gets its own grade record",
            "▸ Admin UI shows 🔄 Retake badge on profile",
        ]),
    ]

    for i, (title, color, bullets) in enumerate(retake_data):
        x = Inches(0.35 + i * 3.15)
        y = Inches(1.7)
        card = add_card(slide, x, y, Inches(3.0), Inches(2.8))

        add_accent_bar(slide, x + Inches(0.2), y + Inches(0.15),
                       Inches(2.6), Inches(0.05), color)

        add_text_box(slide, x + Inches(0.2), y + Inches(0.3),
                     Inches(2.6), Inches(0.3),
                     title, font_size=14, bold=True, color=color)

        add_bullet_slide_content(slide, bullets,
                                  x + Inches(0.2), y + Inches(0.7),
                                  Inches(2.6), Inches(2.0),
                                  font_size=10, color=LIGHT_GRAY, spacing=Pt(6))


def slide_google_sheets(prs):
    """Slide 12 — Google Sheets Sync."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    add_section_header(slide, "Google Sheets Sync",
                       "Automatic cloud backup of every registration")

    items = [
        "▸  Every new student registration automatically appends a row to a configured Google Sheet",
        "▸  Uses Google Sheets API v4 via a Service Account with Editor access",
        "▸  Synced fields: Student ID, name, email, phone, gender, DOB, batch code, and more",
        "▸  Supabase remains the source of truth — Sheets is a complementary backup",
        "▸  If Google Sheets env vars are not configured, the app logs a warning and continues",
        "▸  Configurable sheet tab name (default: 'Registrations')",
        "▸  No data loss — sync failures are caught and logged without blocking the registration flow",
    ]
    add_bullet_slide_content(slide, items,
                              Inches(0.6), Inches(1.7), Inches(8.8), Inches(3.5),
                              font_size=14, color=WHITE, spacing=Pt(12))


def slide_security(prs):
    """Slide 13 — Security & Authentication."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    add_section_header(slide, "Security & Authentication",
                       "Layered security model protecting admin and student data")

    items = [
        "🔐  Admin Authentication — Supabase email/password auth via client SDK",
        "🛡️  Admin Guard Hook — useAdminGuard() checks session on every admin page load",
        "🔑  Service Role Key — Server-side API routes use the Supabase service_role key (bypasses RLS)",
        "📋  Row-Level Security (RLS) — Enforced on all tables:",
        "        • batches: public can only READ active batches",
        "        • students: NO public access; admin-only reads",
        "        • student_grades / mit_grades / proclaimers_grades: admin-only",
        "📸  Storage Policies — student-photos bucket allows public uploads and reads",
        "🔒  Bearer Token Auth — All API routes validate the Authorization header",
        "🚫  Cascade Protection — Foreign keys use ON DELETE RESTRICT to prevent orphaned records",
    ]
    add_bullet_slide_content(slide, items,
                              Inches(0.6), Inches(1.7), Inches(8.8), Inches(3.8),
                              font_size=12.5, color=WHITE, spacing=Pt(8))


def slide_registration_form(prs):
    """Slide 14 — Registration Form Details."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    add_section_header(slide, "Registration Form",
                       "What students see when they open the registration link")

    left_items = [
        "Personal Details:",
        "▸ Surname, First Name, Middle Name",
        "▸ Email Address",
        "▸ Phone Number (WhatsApp)",
        "▸ Date of Birth",
        "▸ Gender (Male / Female)",
        "▸ First Timer status (Yes / No)",
        "▸ Home Address",
        "▸ Next of Kin + Address",
        "▸ State of Origin & Nationality",
        "▸ Education Level (3 tiers)",
    ]

    right_items = [
        "Spiritual Background:",
        "▸ Born Again (Yes / No / Maybe)",
        "    → If Yes: when & where details",
        "▸ Baptized by Immersion (Water)?",
        "    → If Yes: when & where details",
        "▸ Baptized in the Holy Spirit?",
        "    → If Yes: when & where details",
        "▸ When did you join CGCC?",
        "▸ Any challenges to participation?",
        "",
        "▸ Optional: Photo Upload",
    ]

    add_bullet_slide_content(slide, left_items,
                              Inches(0.4), Inches(1.7), Inches(4.5), Inches(3.8),
                              font_size=12, color=LIGHT_GRAY, spacing=Pt(5))
    add_bullet_slide_content(slide, right_items,
                              Inches(5.0), Inches(1.7), Inches(4.5), Inches(3.8),
                              font_size=12, color=LIGHT_GRAY, spacing=Pt(5))


def slide_database_schema(prs):
    """Slide 15 — Database Schema Overview."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    add_section_header(slide, "Database Schema",
                       "6 core tables powering the application")

    tables = [
        ("batches", "id, batch_code, batch_name,\nreg_token, is_active, programme_type", BLUE),
        ("students", "id, student_unique_id, batch_id,\nbiodata fields, photo_url, card_number", GREEN),
        ("student_grades", "student_id → students,\n6 score columns + status + comments", GOLD),
        ("mit_registrations", "batch_id, membership_student_id,\ndepartment (links Membership → MIT)", AMBER),
        ("mit_grades", "mit_registration_id,\n10+ score columns + status", PURPLE),
        ("proclaimers_*", "Same pattern as MIT:\nregistration table + grades table", RED),
    ]

    for i, (name, desc, color) in enumerate(tables):
        col = i % 3
        row = i // 3
        x = Inches(0.35 + col * 3.15)
        y = Inches(1.7 + row * 1.7)
        card = add_card(slide, x, y, Inches(3.0), Inches(1.45))

        add_accent_bar(slide, x + Inches(0.15), y + Inches(0.12),
                       Inches(0.4), Inches(0.04), color)

        add_text_box(slide, x + Inches(0.15), y + Inches(0.22),
                     Inches(2.7), Inches(0.3),
                     name, font_size=13, bold=True, color=color)

        add_text_box(slide, x + Inches(0.15), y + Inches(0.55),
                     Inches(2.7), Inches(0.85),
                     desc, font_size=10, color=MID_GRAY)


def slide_deployment(prs):
    """Slide 16 — Deployment & Next Steps."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    add_section_header(slide, "Deployment & Next Steps")

    items_left = [
        "Deployment:",
        "▸ Next.js 14 deployed on Vercel",
        "▸ Environment variables configured via",
        "   Vercel Project Settings",
        "▸ Supabase hosted (cloud PostgreSQL)",
        "▸ 0.0.0.0 binding for local network testing",
        "▸ Production URL set via NEXT_PUBLIC_BASE_URL",
    ]

    items_right = [
        "Potential Enhancements:",
        "▸ Supabase SSR middleware for edge-level",
        "   admin route protection",
        "▸ Custom Student ID format patterns",
        "▸ Batch activation/deactivation toggle in UI",
        "▸ Advanced analytics: pass/fail rates,",
        "   cohort comparisons, grade distributions",
        "▸ Email/SMS notifications on registration",
        "▸ Bulk certificate generation",
    ]

    add_bullet_slide_content(slide, items_left,
                              Inches(0.4), Inches(1.6), Inches(4.5), Inches(3.5),
                              font_size=13, color=LIGHT_GRAY, spacing=Pt(6))
    add_bullet_slide_content(slide, items_right,
                              Inches(5.0), Inches(1.6), Inches(4.5), Inches(3.5),
                              font_size=13, color=LIGHT_GRAY, spacing=Pt(6))


def slide_thank_you(prs):
    """Slide 17 — Thank You / Closing."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_NAVY)

    add_accent_bar(slide, Inches(3.2), Inches(1.8), Inches(3.6), Inches(0.05))

    add_text_box(slide, Inches(1.0), Inches(2.0), Inches(8.0), Inches(0.8),
                 "Thank You", font_size=42, bold=True, color=GOLD,
                 alignment=PP_ALIGN.CENTER)

    add_accent_bar(slide, Inches(3.2), Inches(2.8), Inches(3.6), Inches(0.05))

    add_text_box(slide, Inches(1.0), Inches(3.1), Inches(8.0), Inches(0.5),
                 "ATS Membership App — Citadel Global Community Church",
                 font_size=14, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1.0), Inches(3.65), Inches(8.0), Inches(0.4),
                 "Questions?  •  Feedback?  •  Feature Requests?",
                 font_size=13, color=WHITE, alignment=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════════
# MAIN
# ═════════════════════════════════════════════════════════════════

def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)  # 16:9

    slide_title(prs)
    slide_agenda(prs)
    slide_overview(prs)
    slide_architecture(prs)
    slide_workflow(prs)
    slide_three_programmes(prs)
    slide_admin_dashboard(prs)
    slide_batch_management(prs)
    slide_student_profiles(prs)
    slide_registration_form(prs)
    slide_import_export(prs)
    slide_retake(prs)
    slide_google_sheets(prs)
    slide_security(prs)
    slide_database_schema(prs)
    slide_deployment(prs)
    slide_thank_you(prs)

    out_path = os.path.join(os.path.dirname(__file__), "ATS_Membership_App_Presentation.pptx")
    prs.save(out_path)
    print(f"[OK] Presentation saved to: {out_path}")
    print(f"     Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
